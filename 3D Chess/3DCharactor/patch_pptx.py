"""
將系統架構文字描述填入計畫簡報檔.pptx 的第 6、7、8 頁。
- 內容框高 3.26"，可用行數約 9–10 行（14pt + spcBef 6pt）
- 第 6 頁右側有圖片，文字保持簡短
- 排版沿用原有 XML 格式（spcBef=600、lang=zh-TW、noAutofit）
"""
import os
from lxml import etree
from pptx import Presentation

SRC = os.path.join(os.path.dirname(os.path.abspath(__file__)), '計畫簡報檔.pptx')
DST = os.path.join(os.path.dirname(os.path.abspath(__file__)), '計畫簡報檔_updated.pptx')

NS  = 'http://schemas.openxmlformats.org/drawingml/2006/main'
IND = 228600   # 0.25 inch in EMU

# ── XML 段落建構器 ──────────────────────────────────────────────────────────────
def make_para(text, bold=False, spc=600, indent=0, font_size_pt=None):
    p = etree.Element(f'{{{NS}}}p')

    pPr = etree.SubElement(p, f'{{{NS}}}pPr',
        indent=str(indent), lvl='0', marL=str(indent), rtl='0', algn='l')
    sb = etree.SubElement(pPr, f'{{{NS}}}spcBef')
    etree.SubElement(sb, f'{{{NS}}}spcPts', val=str(spc))
    sa = etree.SubElement(pPr, f'{{{NS}}}spcAft')
    etree.SubElement(sa, f'{{{NS}}}spcPts', val='0')
    etree.SubElement(pPr, f'{{{NS}}}buNone')

    if text:
        r = etree.SubElement(p, f'{{{NS}}}r')
        rPr_attr = {'lang': 'zh-TW'}
        if bold:
            rPr_attr['b'] = '1'
        if font_size_pt:
            rPr_attr['sz'] = str(int(font_size_pt * 100))  # hundredths of a pt
        etree.SubElement(r, f'{{{NS}}}rPr', **rPr_attr)
        t = etree.SubElement(r, f'{{{NS}}}t')
        t.text = text

    etree.SubElement(p, f'{{{NS}}}endParaRPr')
    return p


def set_body(shape, paragraphs):
    """paragraphs: list of (text, bold, indent_emu, font_size_pt|None)"""
    txBody = shape.text_frame._txBody
    for old_p in txBody.findall(f'{{{NS}}}p'):
        txBody.remove(old_p)
    for text, bold, indent, fs in paragraphs:
        txBody.append(make_para(text, bold=bold, indent=indent, font_size_pt=fs))


# ── 各頁內容 ───────────────────────────────────────────────────────────────────
# 格式：(文字, 粗體, 縮排EMU, 字型pt|None)
# 第 6 頁右側有圖片，左側可視寬僅約 1.74"，每行盡量短
# 目標 9 行以內（含空行）

SLIDE6 = [
    # 9 行
    ('三層式架構',                  True,  0,   None),
    ('表示層',                      True,  0,   None),
    ('HTML / CSS / JS / WebSocket', False, IND, None),
    ('業務邏輯層',                  True,  0,   None),
    ('遊戲邏輯・影像辨識・身份驗證', False, IND, None),
    ('通訊：Socket.IO / WebSocket',  False, IND, None),
    ('資料層',                      True,  0,   None),
    ('MySQL・Redis',                False, IND, None),
    ('AWS S3 / MinIO',              False, IND, None),
]

# 第 7 頁：全寬，9 行以內
SLIDE7 = [
    ('客戶端模組',                           True,  0,   None),
    ('登入大廳・遊戲棋盤・題目作答・遊戲結算', False, IND, None),
    ('伺服器端模組',                          True,  0,   None),
    ('使用者管理：帳號驗證、JWT 權限控管',     False, IND, None),
    ('遊戲邏輯：回合、移動、金錢計算、評分',   False, IND, None),
    ('影像辨識：手勢、骰子偵測、ArUco 標記',  False, IND, None),
    ('後台管理',                              True,  0,   None),
    ('題目 CMS・遊戲數據監控・帳號管理',       False, IND, None),
]

# 第 8 頁：全寬，9 行以內
SLIDE8 = [
    ('使用技術',                              True,  0,   None),
    ('前端：HTML5 / JS / WebSocket / Three.js', False, IND, None),
    ('後端：FastAPI / Node.js / OpenCV',       False, IND, None),
    ('資料庫：MySQL / PostgreSQL・Redis',       False, IND, None),
    ('API 設計',                               True,  0,   None),
    ('/api/auth  登入・註冊',                  False, IND, None),
    ('/api/game  開始・移動・題目・排行榜',    False, IND, None),
    ('部署',                                   True,  0,   None),
    ('CloudFlare > Nginx > FastAPI + Node.js', False, IND, None),
]

# ── 套用並儲存 ─────────────────────────────────────────────────────────────────
prs = Presentation(SRC)

for slide_idx, content in [(5, SLIDE6), (6, SLIDE7), (7, SLIDE8)]:
    slide = prs.slides[slide_idx]
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    if len(text_shapes) >= 2:
        set_body(text_shapes[1], content)
        print(f'Page {slide_idx+1} OK ({len(content)} lines)')
    else:
        print(f'Page {slide_idx+1} SKIP')

prs.save(DST)
print(f'\nSaved: {DST}')
