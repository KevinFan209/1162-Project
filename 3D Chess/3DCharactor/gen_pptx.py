from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = Inches(13.33)
H = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 完全空白版面

BLACK  = RGBColor(0x00, 0x00, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
GRAY   = RGBColor(0x55, 0x55, 0x55)

FONT = 'Microsoft JhengHei'

# ── Helper：加文字框 ──────────────────────────────────────────────────────────
def tb(slide, text, x, y, w, h, size, bold=False, color=BLACK,
       align=PP_ALIGN.LEFT, wrap=True):
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return shape

def row(tf, text, size, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, indent=0):
    """在現有 text_frame 新增一段落"""
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return p

def section(slide, title, x, y, w, h, items, title_size=18, item_size=13,
            title_color=DARK):
    """標題 + 條列清單的區塊"""
    shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    tf.word_wrap = True
    # 標題行
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(title_size)
    r0.font.bold = True
    r0.font.color.rgb = title_color
    r0.font.name = FONT
    # 條列
    for item in items:
        row(tf, item, item_size, color=BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# 第 1 頁　系統概述 ＆ 三層架構
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# 頁首標題
tb(s1, 'PyPoly 系統架構', Inches(0.4), Inches(0.2),
   Inches(12.5), Inches(0.7), 36, bold=True, color=DARK,
   align=PP_ALIGN.CENTER)

# 副標題
tb(s1, '遊戲化 Python 學習平台　|　1162 資管系專題',
   Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.4), 16,
   color=GRAY, align=PP_ALIGN.CENTER)

# 分隔線（純文字 em-dash 行）
tb(s1, '─' * 80, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.25),
   10, color=GRAY, align=PP_ALIGN.CENTER)

# ── 左欄：系統概述 ────────────────────────────────────────────────────────────
section(s1, '系統概述', Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.5),
    ['本系統以大富翁遊戲形式結合影像辨識，讓學習者在遊戲過程中',
     '透過手勢辨識、骰子偵測與程式題目作答，達到寓教於樂的效果。',
     '',
     '目標',
     '  ● 互動式 Python 程式語言學習環境',
     '  ● 影像辨識技術增加遊戲互動性',
     '  ● 金錢累積制遊戲模式提升學習動機',
     '  ● 支援線上 / 線下雙模式遊玩',
     '  ● 完整後台管理系統（題庫、數據、帳號）',
     '',
     '目標客群',
     '  ● 資訊管理相關科系學生',
     '  ● 對程式設計有興趣者',
     '  ● 相關教育機構'],
    title_size=20, item_size=13)

# ── 右欄：三層架構 ────────────────────────────────────────────────────────────
section(s1, '三層式架構（Three-Tier Architecture）',
   Inches(6.7), Inches(1.6), Inches(6.2), Inches(5.5),
    ['',
     '【表示層　Presentation Layer】',
     '  HTML5 / CSS / JavaScript',
     '  WebSocket 即時通訊  ·  Three.js 3D 棋盤',
     '  登入大廳、遊戲棋盤、題目作答、結算、後台',
     '',
     '【業務邏輯層　Business Logic Layer】',
     '  遊戲邏輯　→  FastAPI / Node.js',
     '  影像辨識　→  Python + OpenCV（手勢 / 骰子 / ArUco）',
     '  使用者管理 →  JWT + bcrypt',
     '  後台 CMS　→  題庫維護 / 遊戲監控 / 帳號管理',
     '',
     '【資料層　Data Layer】',
     '  MySQL / PostgreSQL　持久化資料',
     '  Redis　　　　　　　 即時遊戲狀態快取',
     '  AWS S3 / MinIO　　　靜態資源儲存'],
    title_size=20, item_size=13)

# 頁碼
tb(s1, '1 / 3', Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.3),
   11, color=GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 第 2 頁　模組設計 ＆ 資料庫設計
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)

tb(s2, '模組設計　＆　資料庫設計',
   Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7),
   36, bold=True, color=DARK, align=PP_ALIGN.CENTER)

tb(s2, '─' * 80, Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.25),
   10, color=GRAY, align=PP_ALIGN.CENTER)

# ── 左欄：模組設計 ────────────────────────────────────────────────────────────
section(s2, '客戶端模組', Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.5),
    ['  ● 登入與大廳　使用者註冊、登入、遊戲等候室',
     '  ● 主遊戲棋盤　地圖渲染、玩家移動、事件觸發',
     '  ● 互動任務　　題目顯示、作答介面、即時回饋',
     '  ● 遊戲結算　　分數計算、排名顯示、歷史記錄'],
    title_size=18, item_size=13)

section(s2, '伺服器端模組', Inches(0.5), Inches(3.55), Inches(5.8), Inches(3.0),
    ['  ● 使用者管理　帳號驗證、JWT 權限管控',
     '  ● 遊戲邏輯　　回合管理、移動判斷、金錢計算',
     '  ● 題目管理　　題庫維護、難度分級、自動評分',
     '  ● 影像辨識　　手勢偵測、骰子偵測、ArUco 標記',
     '',
     '後台管理模組',
     '  ● Question CMS（題目管理系統）',
     '  ● Vision Debugger（影像辨識偵錯）',
     '  ● Game Analytics（遊戲數據監控）',
     '  ● User Management（玩家帳號管理）'],
    title_size=18, item_size=13)

# ── 右欄：資料庫設計 ──────────────────────────────────────────────────────────
section(s2, '資料庫設計', Inches(6.7), Inches(1.2), Inches(6.2), Inches(5.8),
    ['',
     'users（使用者表）',
     '  user_id (PK)  ·  username  ·  password_hash',
     '  email  ·  created_at  ·  last_login',
     '',
     'game_records（遊戲記錄表）',
     '  game_id (PK)  ·  user_id (FK)  ·  score',
     '  money  ·  rounds  ·  start_time  ·  end_time',
     '',
     'questions（題目表）',
     '  question_id (PK)  ·  category  ·  difficulty (1–5)',
     '  content  ·  answer  ·  explanation',
     '',
     '快取 ＆ 儲存',
     '  Redis　　　→  即時遊戲狀態、Session 快取',
     '  S3 / MinIO →  靜態資源、3D 模型、音效'],
    title_size=20, item_size=13)

tb(s2, '2 / 3', Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.3),
   11, color=GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 第 3 頁　API 設計 ＆ 部署 ＆ 安全性
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)

tb(s3, 'API 設計　＆　部署架構　＆　安全性',
   Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7),
   36, bold=True, color=DARK, align=PP_ALIGN.CENTER)

tb(s3, '─' * 80, Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.25),
   10, color=GRAY, align=PP_ALIGN.CENTER)

# ── 左欄：API ─────────────────────────────────────────────────────────────────
section(s3, 'RESTful API', Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.8),
    ['  POST   /api/auth/login        使用者登入',
     '  POST   /api/auth/register     使用者註冊',
     '  GET    /api/game/:id          取得遊戲資訊',
     '  POST   /api/game/start        開始新遊戲',
     '  POST   /api/game/move         玩家移動',
     '  GET    /api/questions/random  取得隨機題目',
     '  POST   /api/answers/submit    提交答案',
     '  GET    /api/leaderboard       排行榜'],
    title_size=18, item_size=12)

section(s3, 'WebSocket / Socket.IO 事件', Inches(0.5), Inches(3.85), Inches(5.8), Inches(2.8),
    ['  dice_roll      Client → Server   骰子結果傳送',
     '  player_move    Server → Client   玩家位置更新',
     '  question_show  Server → Client   顯示題目',
     '  answer_result  Server → Client   答題結果回傳',
     '  game_state     Server → Client   整體狀態更新'],
    title_size=18, item_size=12)

# ── 右欄：部署 ＆ 安全 ────────────────────────────────────────────────────────
section(s3, '部署架構', Inches(6.7), Inches(1.2), Inches(6.2), Inches(3.5),
    ['生產環境元件',
     '  CloudFlare　→  CDN 加速 + DDoS 防護',
     '  Nginx　　　 →  反向代理 / 靜態資源服務',
     '  Node.js PM2 →  遊戲邏輯 API 程序管理',
     '  FastAPI + Uvicorn  →  影像辨識 ASGI 伺服器',
     '  MySQL　　　 →  持久化關聯資料',
     '  Redis　　　 →  快取伺服器',
     '  AWS S3 / MinIO  →  物件儲存',
     '',
     '開發工具',
     '  VS Code  ·  Node.js  ·  Python + OpenCV  ·  Git'],
    title_size=18, item_size=12)

section(s3, '安全性設計', Inches(6.7), Inches(4.85), Inches(6.2), Inches(2.3),
    ['  ● JWT　　→  無狀態 Token 驗證',
     '  ● bcrypt →  密碼雜湊加密儲存',
     '  ● HTTPS / WSS  →  全程加密通訊',
     '  ● ORM 參數化查詢  →  防 SQL Injection',
     '  ● 輸出轉義 + CSP  →  防 XSS 攻擊'],
    title_size=18, item_size=12)

tb(s3, '3 / 3', Inches(12.5), Inches(7.15), Inches(0.7), Inches(0.3),
   11, color=GRAY, align=PP_ALIGN.RIGHT)

# ── 儲存 ──────────────────────────────────────────────────────────────────────
import os
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'PyPoly_slides.pptx')
prs.save(OUT)
print('Saved:', OUT)
